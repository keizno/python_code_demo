import cv2
import os
import tkinter as tk
import json
from tkinter import ttk, filedialog, messagebox, scrolledtext
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import threading
from pptx.oxml.ns import qn
from lxml import etree
import re
import random

class VideoToPPTXApp:
    def __init__(self, root):
        self.root = root
        self.root.title("🎬 Video to PPTX - v1.2 by sungkb04@khnp.co.kr")
        self.root.geometry("600x900")
        self.root.resizable(False, False)
        self.root.wm_attributes("-topmost", 1)
        
        # 디자인 테마
        self.bg_color = "#F0F4F8"
        self.card_bg = "#FFFFFF"
        self.accent_color = "#4A90E2"
        self.text_color = "#2C3E50"
        self.root.configure(bg=self.bg_color)
        
        # 변수
        self.video_path = tk.StringVar()
        self.times_text = tk.StringVar(value="")
        self.title_text = tk.StringVar(value="Video Presentation")
        self.note_text = tk.StringVar()
        self.existing_ppt_path = tk.StringVar()
        self.add_to_existing = tk.BooleanVar(value=True)
        self.remove_borders = tk.BooleanVar(value=False)
        self.include_audio = tk.BooleanVar(value=False)
        
        # 경로 기억 변수 (독립적 유지)
        self.last_video_dir = os.getcwd()
        self.last_ppt_dir = os.getcwd()
        
        # 설정 파일 경로 (사용자 홈 디렉토리)
        self.config_file = os.path.join(os.path.expanduser("~"), "video_to_pptx_config.json")
        self.load_config()
        
        self.setup_ui()
        self.root.protocol("WM_DELETE_WINDOW", self.on_close)
        
    def setup_ui(self):
        main_frame = tk.Frame(self.root, bg=self.bg_color, padx=20, pady=20)
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        tk.Label(main_frame, text="🎬 PPT Frame Extractor", font=("Segoe UI", 16, "bold"), 
                 bg=self.bg_color, fg=self.text_color).pack(pady=(0, 15))
        
        # 섹션 구성
        self.create_section(main_frame, "📁 비디오 파일", self.create_video_input)
        self.create_section(main_frame, "📎 기존 PPT 업데이트", self.create_ppt_input)
        self.create_section(main_frame, "📝 프레젠테이션 제목 (새 파일용)", self.create_title_input)
        self.create_section(main_frame, "⏱️ 추출 시간 (예: 00:01:10; 46:31,274; 00:02:40-00:02:50)", self.create_time_input)
        self.create_section(main_frame, "🗒️ 슬라이드 노트 메모 (선택 사항)", self.create_note_input)
        self.create_section(main_frame, "⚙️ 옵션", self.create_options)
        self.create_section(main_frame, "🛠️ 도구 (기존 파일 수정)", self.create_tools)
        
        # 실행 버튼
        self.gen_btn = tk.Button(main_frame, text="PPT 생성 및 업데이트 시작", command=self.start_thread, 
                                 bg="#27AE60", fg="white", font=("Segoe UI", 11, "bold"), pady=10, relief=tk.FLAT)
        self.gen_btn.pack(fill=tk.X, pady=(20, 5))
        
        # 프로그레스바 스타일 설정 (노란색)
        style = ttk.Style()
        style.theme_use('clam')
        style.configure("Yellow.Horizontal.TProgressbar", background='#F1C40F', troughcolor='#FFFFFF', bordercolor='#E0E0E0')
        self.progress = ttk.Progressbar(main_frame, style="Yellow.Horizontal.TProgressbar", mode='determinate')
        self.progress.pack(fill=tk.X, pady=5)
        
        # 로그창 추가
        tk.Label(main_frame, text="📜 처리 로그", font=("Segoe UI", 9, "bold"), bg=self.bg_color, fg=self.text_color).pack(anchor="w", pady=(5,0))
        self.log_text = scrolledtext.ScrolledText(main_frame, height=12, state=tk.DISABLED, font=("Consolas", 9))
        self.log_text.pack(fill=tk.BOTH, expand=True, pady=(0, 5))
        
        # 초기 상태 UI 갱신 (설정 로드 후 버튼 상태 동기화)
        self.toggle_ppt()
        if self.video_path.get():
            self.play_btn.config(state=tk.NORMAL)

    def create_section(self, parent, label, widget_func):
        frame = tk.Frame(parent, bg=self.card_bg, highlightbackground="#E1E8ED", highlightthickness=1)
        frame.pack(fill=tk.X, pady=5)
        tk.Label(frame, text=label, font=("Segoe UI", 9, "bold"), bg=self.card_bg, fg=self.text_color).pack(anchor="w", padx=10, pady=(5, 2))
        inner = tk.Frame(frame, bg=self.card_bg)
        inner.pack(fill=tk.X, padx=10, pady=(0, 8))
        widget_func(inner)

    def create_video_input(self, p):
        tk.Entry(p, textvariable=self.video_path, font=("Segoe UI", 9), bg="#F8F9FA").pack(side=tk.LEFT, fill=tk.X, expand=True, ipady=3)
        # 재생 버튼 (가장 오른쪽)
        self.play_btn = tk.Button(p, text="▶ 재생", command=self.play_video, state=tk.DISABLED, bg="#34495E", fg="white", relief=tk.FLAT)
        self.play_btn.pack(side=tk.RIGHT, padx=(5,0))
        # 찾기 버튼 (재생 버튼 왼쪽)
        tk.Button(p, text="찾기", command=self.browse_video, bg=self.accent_color, fg="white", relief=tk.FLAT).pack(side=tk.RIGHT, padx=(5,0))

    def create_ppt_input(self, p):
        tk.Checkbutton(p, text="기존 파일에 슬라이드 추가", variable=self.add_to_existing, command=self.toggle_ppt, bg=self.card_bg).pack(anchor="w")
        self.ppt_ent = tk.Entry(p, textvariable=self.existing_ppt_path, state=tk.DISABLED, bg="#F8F9FA")
        self.ppt_ent.pack(side=tk.LEFT, fill=tk.X, expand=True, ipady=3)
        # 열기 버튼 (가장 오른쪽)
        self.open_ppt_btn = tk.Button(p, text="📂 열기", command=self.open_existing_ppt, state=tk.DISABLED, bg="#34495E", fg="white", relief=tk.FLAT)
        self.open_ppt_btn.pack(side=tk.RIGHT, padx=(5,0))
        # 선택 버튼 (열기 버튼 왼쪽)
        self.ppt_btn = tk.Button(p, text="선택", command=self.browse_ppt, state=tk.DISABLED, bg=self.accent_color, fg="white", relief=tk.FLAT)
        self.ppt_btn.pack(side=tk.RIGHT, padx=(5,0))

    def create_title_input(self, p):
        tk.Entry(p, textvariable=self.title_text, font=("Segoe UI", 9), bg="#F8F9FA").pack(fill=tk.X, ipady=3)

    def create_time_input(self, p):
        tk.Entry(p, textvariable=self.times_text, font=("Segoe UI", 9), bg="#F8F9FA").pack(fill=tk.X, ipady=3)

    def create_note_input(self, p):
        tk.Entry(p, textvariable=self.note_text, font=("Segoe UI", 9), bg="#F8F9FA").pack(fill=tk.X, ipady=3)

    def create_options(self, p):
        tk.Checkbutton(p, text="검은색 테두리 자동 제거 (Auto Crop)", variable=self.remove_borders, bg=self.card_bg).pack(anchor="w")
        tk.Checkbutton(p, text="소리 포함 (moviepy 필요, 속도 느림)", variable=self.include_audio, bg=self.card_bg).pack(anchor="w")

    def create_tools(self, p):
        btn = tk.Button(p, text="선택된 기존 PPT의 모든 비디오에 재생 아이콘 추가하기", command=self.start_add_icons_thread,
                        bg="#95A5A6", fg="white", relief=tk.FLAT)
        btn.pack(fill=tk.X, ipady=2)

    def toggle_ppt(self):
        enabled = self.add_to_existing.get()
        state_common = tk.NORMAL if enabled else tk.DISABLED
        self.ppt_ent.config(state=state_common)
        self.ppt_btn.config(state=state_common)
        
        # 열기 버튼은 활성화 상태이고 경로가 있을 때만 활성화
        if enabled and self.existing_ppt_path.get():
            self.open_ppt_btn.config(state=tk.NORMAL)
        else:
            self.open_ppt_btn.config(state=tk.DISABLED)

    def browse_video(self):
        path = filedialog.askopenfilename(initialdir=self.last_video_dir, filetypes=[("Video", "*.mp4 *.avi *.mov *.mkv")])
        if path: 
            self.video_path.set(path)
            self.last_video_dir = os.path.dirname(path)
            self.play_btn.config(state=tk.NORMAL)

    def browse_ppt(self):
        path = filedialog.askopenfilename(initialdir=self.last_ppt_dir, filetypes=[("PowerPoint", "*.pptx")])
        if path: 
            self.existing_ppt_path.set(path)
            self.last_ppt_dir = os.path.dirname(path)
            self.open_ppt_btn.config(state=tk.NORMAL)

    def load_config(self):
        if os.path.exists(self.config_file):
            try:
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    self.last_video_dir = data.get('last_video_dir', self.last_video_dir)
                    self.last_ppt_dir = data.get('last_ppt_dir', self.last_ppt_dir)
                    self.remove_borders.set(data.get('remove_borders', False))
                    self.include_audio.set(data.get('include_audio', False))
                    self.video_path.set(data.get('video_path', ''))
                    self.existing_ppt_path.set(data.get('existing_ppt_path', ''))
                    self.note_text.set(data.get('note_text', ''))
            except Exception as e:
                print(f"Config load failed: {e}")

    def save_config(self):
        data = {
            'last_video_dir': self.last_video_dir,
            'last_ppt_dir': self.last_ppt_dir,
            'remove_borders': self.remove_borders.get(),
            'include_audio': self.include_audio.get(),
            'video_path': self.video_path.get(),
            'existing_ppt_path': self.existing_ppt_path.get(),
            'note_text': self.note_text.get()
        }
        try:
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=4)
        except Exception:
            pass

    def on_close(self):
        self.save_config()
        self.root.destroy()

    def update_progress(self, value, maximum=None):
        """메인 스레드에서 안전하게 진행률 업데이트"""
        def _update():
            if maximum is not None:
                self.progress.configure(maximum=maximum)
            self.progress.configure(value=value)
        self.root.after(0, _update)

    def start_thread(self):
        if not self.video_path.get(): return messagebox.showerror("알림", "비디오 파일을 선택하세요.")
        if not self.times_text.get().strip(): return messagebox.showwarning("알림", "추출할 시간을 입력하세요.")
        self.gen_btn.config(state=tk.DISABLED)
        # 처음부터 determinate 모드로 초기화 (indeterminate 사용 안 함)
        self.progress.configure(mode='determinate', maximum=100, value=0)
        threading.Thread(target=self.process, daemon=True).start()
        
    def start_add_icons_thread(self):
        if not self.existing_ppt_path.get(): return messagebox.showerror("알림", "기존 PPT 파일을 선택하세요.")
        self.progress.configure(mode='indeterminate')
        self.progress.start(10)
        threading.Thread(target=self.process_existing_icons, daemon=True).start()

    def log(self, msg):
        def _log():
            self.log_text.config(state=tk.NORMAL)
            self.log_text.insert(tk.END, f"{msg}\n")
            self.log_text.see(tk.END)
            self.log_text.config(state=tk.DISABLED)
        self.root.after(0, _log)

    def parse_seconds(self, t_str):
        t_str = t_str.replace(',', '.')
        parts = list(map(float, t_str.split(':')))
        if len(parts) == 3:
            return parts[0]*3600 + parts[1]*60 + parts[2]
        elif len(parts) == 2:
            return parts[0]*60 + parts[1]
        else:
            return parts[0]

    def process(self):
        try:
            is_update = self.add_to_existing.get()
            ppt_target = self.existing_ppt_path.get()
            
            # 1. PPT 로드 또는 생성
            if is_update and ppt_target:
                ppt_target = os.path.normpath(ppt_target)
                if not os.path.exists(ppt_target):
                    raise FileNotFoundError(f"기존 PPT 파일을 찾을 수 없습니다:\n{ppt_target}")
                if not ppt_target.lower().endswith('.pptx'):
                    raise ValueError(f"파일이 .pptx 형식이 아닙니다:\n{ppt_target}")
                try:
                    prs = Presentation(ppt_target)
                except Exception as e:
                    raise IOError(f"PPT 파일을 열 수 없습니다. 파일이 열려 있거나 손상되었을 수 있습니다.\n경로: {ppt_target}\n오류: {e}")
                output_path = ppt_target  # 기존 파일에 덮어쓰기
            else:
                prs = Presentation()
                prs.slide_width = Inches(13.333) # 16:9
                prs.slide_height = Inches(7.5)
                # 새 파일일 때만 제목 슬라이드 추가
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                title_slide.shapes.title.text = self.title_text.get()
                v_name = os.path.splitext(os.path.basename(self.video_path.get()))[0]
                output_path = os.path.join(os.path.dirname(self.video_path.get()), f"{v_name}_slides.pptx")

            cap = cv2.VideoCapture(self.video_path.get())
            if not cap.isOpened():
                raise IOError(f"비디오 파일을 열 수 없습니다:\n{self.video_path.get()}")
            times = [t.strip() for t in self.times_text.get().split(';') if t.strip()]
            
            if not times:
                raise ValueError("추출할 시간을 입력하세요.")

            # 각 아이템당 2단계(추출 + PPT 삽입) → total = len(times) * 2
            total_steps = len(times) * 2
            self.update_progress(0, maximum=total_steps)
            self.log(f"총 {len(times)}개의 작업을 시작합니다. (단계별 진행률 표시)")
            temp_files_to_remove = []

            for idx, t_str in enumerate(times):
                # ── 단계 1: 비디오/이미지 추출 ──────────────────────────
                step_extract = idx * 2 + 1
                self.log(f"[{idx+1}/{len(times)}] 📥 추출 중: {t_str}")
                self.update_progress(step_extract)

                if '-' in t_str:
                    # 비디오 클립 추출 (예: 00:01:10-00:01:20)
                    parts = t_str.split('-', 1)
                    if len(parts) != 2:
                        raise ValueError(f"시간 범위 형식이 잘못되었습니다: '{t_str}'\n예시: 00:01:10-00:01:20")
                    start_str, end_str = parts
                    start_sec = self.parse_seconds(start_str.strip())
                    end_sec = self.parse_seconds(end_str.strip())
                    
                    # 썸네일(포스터) 이미지 추출 (스피커 아이콘 방지)
                    cap.set(cv2.CAP_PROP_POS_MSEC, start_sec * 1000)
                    success, frame = cap.read()
                    poster_path = f"temp_poster_{idx}.png"
                    if success:
                        # 포스터 이미지에 재생 버튼 그리기 (플레이 시 사라짐 효과)
                        frame = self.draw_play_icon_cv2(frame)
                        cv2.imwrite(poster_path, frame)
                    else:
                        poster_path = None

                    tmp_vid = f"temp_clip_{idx}.mp4"
                    self.extract_video_clip(self.video_path.get(), tmp_vid, start_sec, end_sec, self.include_audio.get())

                    # ── 단계 2: PPT 슬라이드 삽입 ─────────────────────────
                    step_insert = idx * 2 + 2
                    self.log(f"[{idx+1}/{len(times)}] 📎 PPT 삽입 중: {t_str}")
                    self.update_progress(step_insert)
                    
                    if os.path.exists(tmp_vid):
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        slide.shapes.add_movie(tmp_vid, 0, 0, width=prs.slide_width, height=prs.slide_height, poster_frame_image=poster_path)
                        
                        # 노트 추가
                        notes_slide = slide.notes_slide
                        note_content = f"Video Clip Path: {self.video_path.get()}\nRange: {t_str}"
                        user_note = self.note_text.get()
                        if user_note:
                            note_content += f"\n\nMemo: {user_note}"
                        notes_slide.notes_text_frame.text = note_content
                        
                        temp_files_to_remove.append(tmp_vid)
                        if poster_path and os.path.exists(poster_path):
                            temp_files_to_remove.append(poster_path)
                else:
                    # 이미지 프레임 추출
                    sec = self.parse_seconds(t_str)
                    cap.set(cv2.CAP_PROP_POS_MSEC, sec * 1000)
                    success, frame = cap.read()
                    
                    if success:
                        # 검은색 테두리 제거 옵션 확인
                        if self.remove_borders.get():
                            frame = self.crop_frame(frame)

                        tmp = f"temp_{idx}.png"
                        cv2.imwrite(tmp, frame)

                        # ── 단계 2: PPT 슬라이드 삽입 ─────────────────────
                        step_insert = idx * 2 + 2
                        self.log(f"[{idx+1}/{len(times)}] 📎 PPT 삽입 중: {t_str}")
                        self.update_progress(step_insert)
                        
                        # 빈 슬라이드(6) 추가 및 이미지 삽입
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        slide.shapes.add_picture(tmp, 0, 0, width=prs.slide_width, height=prs.slide_height)
                        
                        # [유지 기능] 슬라이드 노트에 시간 정보 기록
                        notes_slide = slide.notes_slide
                        note_content = f"Video Frame Path: {self.video_path.get()}\nTimestamp: {t_str}"
                        user_note = self.note_text.get()
                        if user_note:
                            note_content += f"\n\nMemo: {user_note}"
                        notes_slide.notes_text_frame.text = note_content
                        
                        temp_files_to_remove.append(tmp)
                    else:
                        # 추출 실패 시에도 단계2 건너뛰고 진행
                        self.log(f"⚠️ [{idx+1}/{len(times)}] 프레임 추출 실패: {t_str}")
                        self.update_progress(idx * 2 + 2)

            cap.release()
            self.log("💾 PPT 파일 저장 중...")
            prs.save(output_path)
            self.update_progress(total_steps)  # 저장 완료 → 100%
            self.log("✅ 모든 작업이 완료되었습니다.")
            
            # 저장 완료 후 임시 파일 일괄 삭제
            for f in temp_files_to_remove:
                if os.path.exists(f):
                    try: os.remove(f)
                    except: pass

            self.root.after(0, lambda: self.ask_open_file(output_path))
            
        except Exception as e:
            self.root.after(0, lambda err=str(e): messagebox.showerror("오류", err))
        finally:
            self.root.after(0, lambda: self.gen_btn.config(state=tk.NORMAL))

    def draw_play_icon_cv2(self, img):
        """OpenCV를 사용하여 이미지 중앙에 반투명 재생 버튼 그리기"""
        try:
            overlay = img.copy()
            h, w = img.shape[:2]
            center_x, center_y = w // 2, h // 2
            # 크기 계산 (화면의 약 15%)
            radius = int(min(h, w) * 0.08)
            
            # 1. 원 그리기 (검은색 채움 - 배경)
            # 테두리 없이 내부를 검은색으로 채움 (thickness=-1)
            cv2.circle(overlay, (center_x, center_y), radius, (0, 0, 0), -1, lineType=cv2.LINE_AA)
            
            # 2. 삼각형 그리기 (흰색 채움)
            # 삼각형 좌표 계산
            tri_len = int(radius * 0.7)
            pt1 = (center_x + int(tri_len * 1.2), center_y) # 오른쪽 끝
            pt2 = (center_x - int(tri_len * 0.6), center_y - tri_len) # 왼쪽 위
            pt3 = (center_x - int(tri_len * 0.6), center_y + tri_len) # 왼쪽 아래
            triangle_cnt = np.array([pt1, pt2, pt3])
            cv2.drawContours(overlay, [triangle_cnt], 0, (255, 255, 255), -1, lineType=cv2.LINE_AA)
            
            # 3. 투명도 적용 (알파 블렌딩)
            alpha = 0.7  # 불투명도
            cv2.addWeighted(overlay, alpha, img, 1 - alpha, 0, img)
            
        except Exception as e:
            print(f"Icon drawing failed: {e}")
        return img

    def process_existing_icons(self):
        """기존 PPT를 열어 비디오 객체의 포스터 프레임(썸네일)에 재생 아이콘을 합성"""
        try:
            from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
            from pptx.dml.color import RGBColor
            import io
            
            ppt_path = self.existing_ppt_path.get()
            self.log(f"📂 기존 파일 분석 중: {ppt_path}")
            
            prs = Presentation(ppt_path)
            count = 0
            
            for slide in prs.slides:
                # 리스트를 복사해서 순회 (도형 추가/삭제 시 인덱스 문제 방지)
                for shape in list(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                        # 이미 처리된 비디오인지 확인 (중복 방지)
                        if shape.name and " [PlayIcon]" in shape.name:
                            continue

                        try:
                            # 1. 기존 포스터 프레임 이미지 추출
                            # blipFill/blip 요소 접근
                            blip_fill = shape._element.blipFill
                            if blip_fill is None:
                                continue
                                
                            blip = blip_fill.blip
                            if blip is None:
                                continue
                                
                            rId = blip.get(qn('r:embed'))
                            if not rId:
                                continue
                                
                            image_part = slide.part.related_part(rId)
                            image_bytes = image_part.blob
                            
                            # 2. OpenCV로 이미지 로드 및 아이콘 그리기
                            nparr = np.frombuffer(image_bytes, np.uint8)
                            img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                            
                            if img is None:
                                continue
                                
                            # 기존 draw_play_icon_cv2 함수 재사용 (이미지 위에 아이콘 합성)
                            img = self.draw_play_icon_cv2(img)
                            
                            # 3. 이미지를 바이트 스트림으로 변환
                            is_success, buffer = cv2.imencode(".png", img)
                            if not is_success:
                                continue
                            
                            img_stream = io.BytesIO(buffer.tobytes())
                            
                            # 4. 새 이미지를 슬라이드에 임시로 추가하여 등록 (rId 생성)
                            # (python-pptx에서 파트 등록을 위한 가장 안전한 방법)
                            temp_pic = slide.shapes.add_picture(img_stream, 0, 0, width=Inches(1), height=Inches(1))
                            new_rId = temp_pic._element.blipFill.blip.get(qn('r:embed'))
                            
                            # 5. 비디오 쉐이프의 blip을 새 이미지 rId로 교체
                            blip.set(qn('r:embed'), new_rId)
                            
                            # 6. 임시 이미지 쉐이프 삭제
                            # XML 트리에서 직접 제거하여 흔적 남기지 않음
                            temp_pic._element.getparent().remove(temp_pic._element)
                            
                            # 처리 완료 태그 추가 (중복 방지용)
                            shape.name = f"{shape.name} [PlayIcon]"
                            count += 1
                            
                        except Exception as e:
                            print(f"Error processing shape {shape.shape_id}: {e}")
                            continue

            save_path = ppt_path.replace(".pptx", "_icon_added.pptx")
            prs.save(save_path)
            self.log(f"✅ {count}개의 비디오 썸네일 교체 완료.\n(재생 시 아이콘이 자연스럽게 사라지며, 클릭/일시정지 기능이 정상 작동합니다.)")
            self.root.after(0, lambda: self.ask_open_file(save_path))
            
        except Exception as e:
            self.root.after(0, lambda err=str(e): messagebox.showerror("오류", err))
        finally:
            self.root.after(0, self.progress.stop)

    def extract_video_clip(self, input_path, output_path, start_sec, end_sec, include_audio=False):
        # 1. moviepy 시도 (용량 최적화 및 오디오 지원)
        # 소리 미포함 시에도 moviepy(libx264)를 쓰면 용량이 훨씬 작아짐
        try:
            try:
                from moviepy import VideoFileClip
            except ImportError:
                from moviepy.editor import VideoFileClip

            clip = VideoFileClip(input_path)
            try:
                clip = clip.subclipped(start_sec, end_sec)
            except AttributeError:
                clip = clip.subclip(start_sec, end_sec)
            
            # audio=False로 설정하면 소리 없이 영상만 저장 (용량 최적화)
            clip.write_videofile(output_path, codec='libx264', audio_codec='aac', audio=include_audio, logger=None)
            clip.close()
            return
        except ImportError:
            self.log("⚠️ moviepy 미설치. OpenCV로 진행합니다. (용량이 클 수 있음)")
        except Exception as e:
            self.log(f"⚠️ moviepy 오류: {e}. OpenCV로 전환합니다.")

        # 2. OpenCV Fallback (오디오 불가, 용량 큼)
        cap = cv2.VideoCapture(input_path)
        if not cap.isOpened():
            raise IOError(f"비디오 파일을 열 수 없습니다: {input_path}")
        fps = cap.get(cv2.CAP_PROP_FPS)
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        fourcc = cv2.VideoWriter_fourcc(*'mp4v')
        out = cv2.VideoWriter(output_path, fourcc, fps, (width, height))
        
        cap.set(cv2.CAP_PROP_POS_MSEC, start_sec * 1000)
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret or (cap.get(cv2.CAP_PROP_POS_MSEC) / 1000) > end_sec:
                break
            out.write(frame)
        cap.release()
        out.release()

    def crop_frame(self, frame):
        try:
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            # 픽셀 값이 10보다 큰(검은색이 아닌) 영역 찾기
            _, thresh = cv2.threshold(gray, 10, 255, cv2.THRESH_BINARY)
            coords = cv2.findNonZero(thresh)
            if coords is not None:
                x, y, w, h = cv2.boundingRect(coords)
                return frame[y:y+h, x:x+w]
        except Exception:
            pass
        return frame

    def _open_file(self, path):
        """크로스플랫폼 파일 열기"""
        import sys, subprocess
        try:
            if sys.platform == 'win32':
                os.startfile(path)
            elif sys.platform == 'darwin':
                subprocess.Popen(['open', path])
            else:
                subprocess.Popen(['xdg-open', path])
        except Exception as e:
            messagebox.showerror("오류", f"파일을 열 수 없습니다: {e}")

    def play_video(self):
        if self.video_path.get():
            self._open_file(self.video_path.get())

    def open_existing_ppt(self):
        if self.existing_ppt_path.get():
            self._open_file(self.existing_ppt_path.get())

    def ask_open_file(self, path):
        if messagebox.askyesno("완료", f"작업이 완료되었습니다.\n생성된 PPT 파일을 여시겠습니까?\n\n경로: {path}"):
            self._open_file(path)

if __name__ == "__main__":
    root = tk.Tk()
    VideoToPPTXApp(root)
    root.mainloop()